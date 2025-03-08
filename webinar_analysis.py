import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Load the Excel file
file_path = "/home/agniswarmukherjee/Desktop/Webinar_Analysis/Webinar Assignment.xlsx"
data = pd.read_excel(file_path, sheet_name="Sheet1")

# Data Preparation
data['Comments - Detailed'] = data['Comments - Detailed'].fillna("No Comment")
relevant_columns = ["Attended (Y/N)", "Event Name", "Status - Dont add extra dropdown", "Comments - Detailed"]
data = data[relevant_columns]

def categorize_comment(comment):
    if isinstance(comment, str):
        comment = comment.lower()
        if "not interested" in comment:
            return "Lack of Interest"
        elif "price" in comment or "budget" in comment or "afford" in comment:
            return "Price Concerns"
        elif "not a fit" in comment or "not suitable" in comment:
            return "Product Fit"
        elif "busy" in comment or "call back" in comment:
            return "Timing Issues"
        elif "incorrect" in comment or "invalid" in comment:
            return "Incorrect Contact"
        else:
            return "Other"
    else:
        return "Other"

data['Category'] = data['Comments - Detailed'].apply(categorize_comment)
category_counts = data['Category'].value_counts()
attendance_counts = data['Attended (Y/N)'].value_counts()
event_counts = data['Event Name'].value_counts()

# Create a presentation
prs = Presentation()

# Slide 1: Overview of the Webinar and Cohorts
slide_layout = prs.slide_layouts[5]  # Title Only layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Webinar Analysis & Lead Conversion Insights"

content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(2.5))
text_frame = content.text_frame
text_frame.text = f"""
Total Leads: {data.shape[0]}
Attendance Rate: {attendance_counts['Yes']} out of {data.shape[0]}
Top 3 Events:
1. {event_counts.index[0]} - {event_counts.iloc[0]}
2. {event_counts.index[1]} - {event_counts.iloc[1]}
3. {event_counts.index[2]} - {event_counts.iloc[2]}
"""

# Slide 2: Key Insights from Sales Comments
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Key Insights from Sales Comments"

content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(2.5))
text_frame = content.text_frame
text_frame.text = f"""
Top Reasons for Non-Conversion:
1. {category_counts.index[0]} - {category_counts.iloc[0]}
2. {category_counts.index[1]} - {category_counts.iloc[1]}
3. {category_counts.index[2]} - {category_counts.iloc[2]}
"""

# Slide 3: Recommendations for Improvement
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Recommendations for Improvement"

content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(2.5))
text_frame = content.text_frame
text_frame.text = """
1. Address Price Concerns with Flexible Pricing
2. Improve Product Demo to Highlight Fit
3. Enhance Follow-up Strategies for Better Engagement
"""

# Save the presentation
prs.save("/home/agniswarmukherjee/Desktop/Webinar_Analysis/Webinar_Analysis_Presentation.pptx")
print("Presentation saved successfully!")

# Optional: Save Chart as an Image and Insert into Presentation

# Plot category counts and save as an image (if not already done)
plt.figure(figsize=(8, 5))
category_counts.plot(kind='bar', color=['blue', 'red', 'green', 'orange', 'purple'])
plt.title("Reasons for Non-Conversion")
plt.xlabel("Category")
plt.ylabel("Count")
plt.xticks(rotation=45)
plt.tight_layout()
plt.savefig("/home/agniswarmukherjee/Desktop/Webinar_Analysis/non_conversion_chart.png")

# Slide for the Chart
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Reasons for Non-Conversion"

# Insert image into the slide
img_path = "/home/agniswarmukherjee/Desktop/Webinar_Analysis/non_conversion_chart.png"
left = Inches(1)
top = Inches(1.5)
slide.shapes.add_picture(img_path, left, top, width=Inches(6))

# Save the updated presentation with chart included
prs.save("/home/agniswarmukherjee/Desktop/Webinar_Analysis/Webinar_Analysis_Presentation.pptx")
print("Presentation with chart saved successfully!")
